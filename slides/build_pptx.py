#!/usr/bin/env python3
"""Build a wide-screen PPTX from rendered slide images and presenter notes."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


HERE = Path(__file__).resolve().parent
PNG_DIR = HERE / "slides_png"
OUTPUT = HERE.parent / "X_Arcade_deck.pptx"
document = json.loads((HERE / "_deck_data.json").read_text(encoding="utf-8"))
slides = document.get("result", document)["slides"]

presentation = Presentation()
presentation.slide_width = Inches(13.333)
presentation.slide_height = Inches(7.5)
blank = presentation.slide_layouts[6]

for index, slide_data in enumerate(slides, start=1):
    image_path = PNG_DIR / f"slide_{index:02d}.png"
    if not image_path.is_file():
        raise FileNotFoundError(image_path)
    slide = presentation.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(image_path),
        0,
        0,
        width=presentation.slide_width,
        height=presentation.slide_height,
    )
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")

presentation.save(OUTPUT)
print(f"wrote {OUTPUT.name} with {len(slides)} slides")
